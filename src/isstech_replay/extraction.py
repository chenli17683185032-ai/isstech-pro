"""Parse immutable material blobs into versioned structured documents."""

from __future__ import annotations

from dataclasses import asdict
from datetime import UTC, datetime
import hashlib
from importlib.metadata import version
import json
import math
import os
from pathlib import Path
import shutil
from uuid import uuid4

from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from .ai.base import ExtractionProvider
from .ai.provider import ProviderResponseError
from .field_mapping import (
    DEFAULT_CONFIDENCE_THRESHOLD,
    field_profile,
    validate_proposals,
)
from .materials import MaterialService
from .models.extraction import (
    DocumentParseResult,
    DocumentUnit,
    SourceKind,
    StructuredDocument,
    ExtractionResult,
    ExtractionStatus,
)
from .models.materials import Material, MaterialStatus
from .sync import safe_error_message
from .validation import require_path_segment


PARSER_CONTRACT_VERSION = "structured-document/1"
DEFAULT_MAX_UNIT_CHARS = 200_000
DEFAULT_MAX_DOCUMENT_CHARS = 2_000_000


class DocumentExtractionError(RuntimeError):
    """A material cannot be converted into a bounded structured document."""


class UnsupportedDocumentType(DocumentExtractionError):
    """No parser is registered for the detected MIME type."""


def _clean_text(value: str) -> str:
    return value.replace("\x00", "").replace("\r\n", "\n").replace("\r", "\n").strip()


def _package_version(name: str) -> str:
    return version(name)


class DocumentExtractor:
    def __init__(
        self,
        *,
        max_unit_chars: int = DEFAULT_MAX_UNIT_CHARS,
        max_document_chars: int = DEFAULT_MAX_DOCUMENT_CHARS,
    ) -> None:
        if max_unit_chars < 1 or max_document_chars < 1:
            raise ValueError("document character limits must be positive")
        self.max_unit_chars = max_unit_chars
        self.max_document_chars = max_document_chars

    def extract(self, material: Material, path: Path) -> StructuredDocument:
        mime_type = material.detected_mime_type
        if mime_type == "application/pdf":
            parser_name = "pypdf"
            parser_version = _package_version("pypdf")
            units, issues = self._pdf(path)
        elif mime_type.endswith("wordprocessingml.document"):
            parser_name = "python-docx"
            parser_version = _package_version("python-docx")
            units, issues = self._docx(path)
        elif mime_type.endswith("spreadsheetml.sheet"):
            parser_name = "openpyxl"
            parser_version = _package_version("openpyxl")
            units, issues = self._xlsx(path)
        elif mime_type.endswith("presentationml.presentation"):
            parser_name = "python-pptx"
            parser_version = _package_version("python-pptx")
            units, issues = self._pptx(path)
        elif mime_type in {"text/plain", "application/json"}:
            parser_name = "utf8-text"
            parser_version = "1"
            units, issues = self._text(path)
        else:
            raise UnsupportedDocumentType(f"unsupported detected MIME type: {mime_type}")

        bounded_units, limit_issues = self._bounded(units)
        all_issues = list(issues)
        all_issues.extend(limit_issues)
        if material.status is MaterialStatus.NEEDS_REVIEW:
            all_issues.append("material MIME/signature review is unresolved")
        if not any(unit.text for unit in bounded_units):
            all_issues.append("document contains no extractable text; OCR/manual review required")
        return StructuredDocument(
            material_id=material.id,
            material_sha256=material.sha256,
            detected_mime_type=mime_type,
            parser_name=parser_name,
            parser_version=f"{PARSER_CONTRACT_VERSION}+{parser_version}",
            units=tuple(bounded_units),
            issues=tuple(dict.fromkeys(all_issues)),
        )

    @staticmethod
    def _pdf(path: Path) -> tuple[list[DocumentUnit], list[str]]:
        reader = PdfReader(path, strict=False)
        if reader.is_encrypted:
            try:
                unlocked = reader.decrypt("")
            except Exception as exc:
                raise DocumentExtractionError("encrypted PDF cannot be opened") from exc
            if not unlocked:
                raise DocumentExtractionError("encrypted PDF requires a password")
        units: list[DocumentUnit] = []
        issues: list[str] = []
        for index, page in enumerate(reader.pages, start=1):
            try:
                text = _clean_text(page.extract_text() or "")
            except Exception as exc:
                text = ""
                issues.append(f"page {index} text extraction failed: {type(exc).__name__}")
            if not text:
                issues.append(f"page {index} has no text layer")
            units.append(
                DocumentUnit(
                    kind=SourceKind.PAGE,
                    index=index,
                    label=f"Page {index}",
                    text=text,
                )
            )
        return units, issues

    @staticmethod
    def _docx(path: Path) -> tuple[list[DocumentUnit], list[str]]:
        document = DocxDocument(path)
        lines = [_clean_text(paragraph.text) for paragraph in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                lines.append("\t".join(_clean_text(cell.text) for cell in row.cells))
        text = "\n".join(line for line in lines if line)
        return [DocumentUnit(SourceKind.DOCUMENT, 1, "Document", text)], []

    @staticmethod
    def _xlsx(path: Path) -> tuple[list[DocumentUnit], list[str]]:
        units: list[DocumentUnit] = []
        with path.open("rb") as source:
            workbook = load_workbook(
                source,
                read_only=True,
                data_only=True,
                keep_links=False,
            )
            try:
                for index, sheet in enumerate(workbook.worksheets, start=1):
                    lines: list[str] = []
                    for row in sheet.iter_rows(values_only=True):
                        values = ["" if value is None else str(value) for value in row]
                        if any(values):
                            lines.append("\t".join(values).rstrip())
                    units.append(
                        DocumentUnit(
                            kind=SourceKind.SHEET,
                            index=index,
                            label=sheet.title,
                            text=_clean_text("\n".join(lines)),
                        )
                    )
            finally:
                workbook.close()
        return units, []

    @staticmethod
    def _pptx(path: Path) -> tuple[list[DocumentUnit], list[str]]:
        presentation = Presentation(path)
        units: list[DocumentUnit] = []
        for index, slide in enumerate(presentation.slides, start=1):
            lines: list[str] = []
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text = _clean_text(shape.text)
                    if text:
                        lines.append(text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        lines.append("\t".join(_clean_text(cell.text) for cell in row.cells))
            units.append(
                DocumentUnit(
                    kind=SourceKind.SLIDE,
                    index=index,
                    label=f"Slide {index}",
                    text=_clean_text("\n".join(lines)),
                )
            )
        return units, []

    @staticmethod
    def _text(path: Path) -> tuple[list[DocumentUnit], list[str]]:
        try:
            text = path.read_text(encoding="utf-8-sig")
        except UnicodeDecodeError as exc:
            raise DocumentExtractionError("text material is not valid UTF-8") from exc
        return [DocumentUnit(SourceKind.DOCUMENT, 1, "Document", _clean_text(text))], []

    def _bounded(
        self,
        units: list[DocumentUnit],
    ) -> tuple[list[DocumentUnit], list[str]]:
        output: list[DocumentUnit] = []
        issues: list[str] = []
        remaining = self.max_document_chars
        for unit in units:
            allowed = min(self.max_unit_chars, remaining)
            text = unit.text
            if len(text) > allowed:
                text = text[:allowed]
                issues.append(
                    f"{unit.kind.value} {unit.index} text truncated at {allowed} characters"
                )
            output.append(
                DocumentUnit(
                    kind=unit.kind,
                    index=unit.index,
                    label=unit.label,
                    text=text,
                )
            )
            remaining -= len(text)
            if remaining <= 0:
                if len(output) < len(units):
                    issues.append("document text truncated at total character limit")
                break
        return output, issues


class DocumentExtractionService:
    def __init__(
        self,
        material_service: MaterialService,
        extractor: DocumentExtractor | None = None,
    ) -> None:
        self.material_service = material_service
        self.extractor = extractor or DocumentExtractor()

    def parse(self, material_id: str) -> DocumentParseResult:
        material = self.material_service.get(material_id)
        if material is None:
            raise ValueError("material not found")
        original = self.material_service.resolve_original(material)
        document = self.extractor.extract(material, original)
        return self._persist(material, document)

    def _persist(
        self,
        material: Material,
        document: StructuredDocument,
    ) -> DocumentParseResult:
        payload = asdict(document)
        canonical = json.dumps(
            payload,
            ensure_ascii=False,
            separators=(",", ":"),
            sort_keys=True,
        )
        document_hash = hashlib.sha256(canonical.encode("utf-8")).hexdigest()
        material_derived = self.material_service.derived_directory(material.id)
        documents_root = material_derived / "documents"
        final_directory = documents_root / document_hash
        documents_root.mkdir(mode=0o700, parents=True, exist_ok=True)
        os.chmod(documents_root, 0o700)

        staging = material_derived / f".staging-{uuid4().hex}"
        staging.mkdir(mode=0o700, parents=True)
        try:
            document_file = staging / "document.json"
            text_file = staging / "text.txt"
            units_directory = staging / "pages"
            units_directory.mkdir(mode=0o700)
            self._write_text(
                document_file,
                json.dumps(payload, ensure_ascii=False, indent=2, sort_keys=True) + "\n",
            )
            combined = []
            unit_names: list[str] = []
            for unit in document.units:
                name = f"{unit.index:04d}.json"
                unit_names.append(name)
                self._write_text(
                    units_directory / name,
                    json.dumps(asdict(unit), ensure_ascii=False, indent=2, sort_keys=True)
                    + "\n",
                )
                combined.append(
                    f"--- {unit.kind.value} {unit.index}: {unit.label} ---\n{unit.text}"
                )
            self._write_text(text_file, "\n\n".join(combined).rstrip() + "\n")

            if final_directory.exists():
                self._verify_existing_artifacts(staging, final_directory, unit_names)
            else:
                try:
                    os.rename(staging, final_directory)
                except FileExistsError:
                    self._verify_existing_artifacts(staging, final_directory, unit_names)

            document_path = final_directory / "document.json"
            text_path = final_directory / "text.txt"
            unit_paths = tuple(final_directory / "pages" / name for name in unit_names)
            for path in (document_path, text_path, *unit_paths):
                if not path.is_file():
                    raise DocumentExtractionError("derived document artifact is incomplete")
                os.chmod(path, 0o600)
            created_at = datetime.now(UTC).isoformat()
            relative_document = document_path.relative_to(
                self.material_service.data_dir
            ).as_posix()
            relative_text = text_path.relative_to(self.material_service.data_dir).as_posix()
            relative_units = tuple(
                path.relative_to(self.material_service.data_dir).as_posix()
                for path in unit_paths
            )
            self._register_artifact(
                material.id,
                "document_json",
                relative_document,
                document.parser_version,
                document_path,
                created_at,
            )
            self._register_artifact(
                material.id,
                "text",
                relative_text,
                document.parser_version,
                text_path,
                created_at,
            )
            for unit, relative, path in zip(
                document.units,
                relative_units,
                unit_paths,
                strict=True,
            ):
                self._register_artifact(
                    material.id,
                    f"unit:{unit.kind.value}",
                    relative,
                    document.parser_version,
                    path,
                    created_at,
                )
            return DocumentParseResult(
                document=document,
                document_path=relative_document,
                text_path=relative_text,
                unit_paths=relative_units,
            )
        finally:
            if staging.exists():
                shutil.rmtree(staging)

    @staticmethod
    def _write_text(path: Path, content: str) -> None:
        path.write_text(content, encoding="utf-8")
        os.chmod(path, 0o600)

    @staticmethod
    def _verify_existing_artifacts(
        staged_directory: Path,
        final_directory: Path,
        unit_names: list[str],
    ) -> None:
        relative_paths = (
            Path("document.json"),
            Path("text.txt"),
            *(Path("pages") / name for name in unit_names),
        )
        for relative in relative_paths:
            staged = staged_directory / relative
            existing = final_directory / relative
            if (
                existing.is_symlink()
                or not existing.is_file()
                or existing.read_bytes() != staged.read_bytes()
            ):
                raise DocumentExtractionError(
                    "existing derived document does not match deterministic output"
                )

    def _register_artifact(
        self,
        material_id: str,
        kind: str,
        relative_path: str,
        parser_version: str,
        path: Path,
        created_at: str,
    ) -> None:
        data = path.read_bytes()
        self.material_service.storage.register_material_artifact(
            material_id=material_id,
            kind=kind,
            path=relative_path,
            parser_version=parser_version,
            sha256=hashlib.sha256(data).hexdigest(),
            size_bytes=len(data),
            created_at=created_at,
        )


class FieldExtractionService:
    """Run one provider behind evidence gates and persist an auditable result."""

    def __init__(
        self,
        material_service: MaterialService,
        provider: ExtractionProvider,
        document_service: DocumentExtractionService | None = None,
    ) -> None:
        self.material_service = material_service
        self.provider = provider
        self.document_service = document_service or DocumentExtractionService(
            material_service
        )

    def extract(
        self,
        material_id: str,
        *,
        profile: str = "purchase_requisition",
        confidence_threshold: float = DEFAULT_CONFIDENCE_THRESHOLD,
        extraction_id: str | None = None,
    ) -> ExtractionResult:
        if not 0 <= confidence_threshold <= 1:
            raise ValueError("confidence_threshold must be between 0 and 1")
        specs = field_profile(profile)
        if self.material_service.get(material_id) is None:
            raise ValueError("material not found")
        actual_id = require_path_segment(
            extraction_id or uuid4().hex,
            "extraction_id",
        )
        started_at = datetime.now(UTC).isoformat()
        extractor_version = f"field-extraction/1+{self.provider.version}"
        self.material_service.storage.start_extraction(
            extraction_id=actual_id,
            material_id=material_id,
            profile=profile,
            provider=self.provider.name,
            model=self.provider.model,
            extractor_version=extractor_version,
            confidence_threshold=confidence_threshold,
            started_at=started_at,
        )
        running = True
        try:
            parsed = self.document_service.parse(material_id)
            proposals = self.provider.propose(parsed.document, specs)
            if any(
                not isinstance(proposal.confidence, (int, float))
                or not math.isfinite(proposal.confidence)
                for proposal in proposals
            ):
                raise ProviderResponseError("provider returned non-finite confidence")
            validation = validate_proposals(
                parsed.document,
                specs,
                proposals,
                confidence_threshold=confidence_threshold,
            )
            status = (
                ExtractionStatus.SUCCEEDED
                if validation.can_advance
                else ExtractionStatus.NEEDS_REVIEW
            )
            finished_at = datetime.now(UTC).isoformat()
            result_path = (
                Path("materials")
                / "derived"
                / material_id
                / "extractions"
                / actual_id
                / "result.json"
            ).as_posix()
            result = ExtractionResult(
                id=actual_id,
                material_id=material_id,
                provider=self.provider.name,
                model=self.provider.model,
                extractor_version=extractor_version,
                status=status,
                confidence_threshold=confidence_threshold,
                can_advance=validation.can_advance,
                document_path=parsed.document_path,
                result_path=result_path,
                proposals=validation.proposals,
                issues=validation.issues,
                started_at=started_at,
                finished_at=finished_at,
            )
            result_file = self._persist_result(result)
            data = result_file.read_bytes()
            self.material_service.storage.register_material_artifact(
                material_id=material_id,
                kind="extraction_result",
                path=result_path,
                parser_version=extractor_version,
                sha256=hashlib.sha256(data).hexdigest(),
                size_bytes=len(data),
                created_at=finished_at,
            )
            self.material_service.storage.complete_extraction(result, field_specs=specs)
            running = False
            return result
        except Exception as error:
            if running:
                self.material_service.storage.fail_extraction(
                    extraction_id=actual_id,
                    finished_at=datetime.now(UTC).isoformat(),
                    error_type=type(error).__name__,
                    error_message=safe_error_message(error),
                )
            raise

    def _persist_result(self, result: ExtractionResult) -> Path:
        material_derived = self.material_service.derived_directory(result.material_id)
        extractions_root = material_derived / "extractions"
        final_directory = extractions_root / result.id
        staging = material_derived / f".extraction-staging-{uuid4().hex}"
        extractions_root.mkdir(mode=0o700, parents=True, exist_ok=True)
        os.chmod(extractions_root, 0o700)
        staging.mkdir(mode=0o700, parents=True)
        try:
            result_file = staging / "result.json"
            result_file.write_text(
                json.dumps(asdict(result), ensure_ascii=False, indent=2, sort_keys=True)
                + "\n",
                encoding="utf-8",
            )
            os.chmod(result_file, 0o600)
            try:
                os.rename(staging, final_directory)
            except FileExistsError as exc:
                raise DocumentExtractionError(
                    f"extraction output already exists: {result.id}"
                ) from exc
            final_file = final_directory / "result.json"
            os.chmod(final_file, 0o600)
            return final_file
        finally:
            if staging.exists():
                shutil.rmtree(staging)
