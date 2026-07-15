"""Supported documents become bounded, immutable, source-addressable artifacts."""

from __future__ import annotations

from io import BytesIO
import json
from pathlib import Path

from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter
import pytest

from isstech_replay.ai.provider import RuleBasedExtractionProvider
from isstech_replay.extraction import (
    DocumentExtractionError,
    DocumentExtractionService,
    DocumentExtractor,
    FieldExtractionService,
    UnsupportedDocumentType,
)
from isstech_replay.materials import MaterialService
from isstech_replay.models.extraction import ExtractionStatus, SourceKind
from tools import extract_material as extract_cli


def _service(tmp_path: Path) -> MaterialService:
    return MaterialService(data_dir=tmp_path / "data", max_bytes=10 * 1024 * 1024)


def _docx_bytes() -> bytes:
    output = BytesIO()
    document = DocxDocument()
    document.add_heading("REDACTED PROJECT", level=1)
    document.add_paragraph("项目编号：PRJ-001")
    table = document.add_table(rows=1, cols=2)
    table.cell(0, 0).text = "采购方式"
    table.cell(0, 1).text = "公开询价"
    document.save(output)
    return output.getvalue()


def _xlsx_bytes() -> bytes:
    output = BytesIO()
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Project Sheet"
    sheet.append(["项目编号", "PRJ-001"])
    sheet.append(["项目名称", "REDACTED PROJECT"])
    workbook.save(output)
    return output.getvalue()


def _pptx_bytes() -> bytes:
    output = BytesIO()
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "REDACTED PROJECT"
    slide.placeholders[1].text = "项目编号：PRJ-001"
    presentation.save(output)
    return output.getvalue()


def _blank_pdf_bytes() -> bytes:
    output = BytesIO()
    writer = PdfWriter()
    writer.add_blank_page(width=200, height=200)
    writer.write(output)
    return output.getvalue()


def test_text_parse_persists_versioned_artifacts_idempotently(tmp_path: Path) -> None:
    materials = _service(tmp_path)
    ingested = materials.ingest_stream(
        BytesIO("项目编号：PRJ-001\n项目名称：REDACTED PROJECT".encode()),
        original_name="project.txt",
        declared_mime_type="text/plain",
    )
    service = DocumentExtractionService(materials)
    first = service.parse(ingested.material.id)
    second = service.parse(ingested.material.id)

    assert first.document.units[0].kind is SourceKind.DOCUMENT
    assert first.document.units[0].index == 1
    assert "项目编号" in first.document.units[0].text
    assert second.document_path == first.document_path
    assert materials.storage.table_count("material_artifacts") == 3

    document_path = materials.data_dir / first.document_path
    text_path = materials.data_dir / first.text_path
    unit_path = materials.data_dir / first.unit_paths[0]
    assert document_path.stat().st_mode & 0o777 == 0o600
    assert text_path.stat().st_mode & 0o777 == 0o600
    assert unit_path.stat().st_mode & 0o777 == 0o600
    assert json.loads(document_path.read_text())["material_id"] == ingested.material.id
    assert materials.resolve_original(ingested.material).stat().st_mode & 0o777 == 0o400


def test_existing_corrupt_derived_artifact_is_not_silently_reused(tmp_path: Path) -> None:
    materials = _service(tmp_path)
    ingested = materials.ingest_stream(
        BytesIO("项目编号：PRJ-001".encode()),
        original_name="project.txt",
        declared_mime_type="text/plain",
    )
    service = DocumentExtractionService(materials)
    parsed = service.parse(ingested.material.id)
    document_path = materials.data_dir / parsed.document_path
    document_path.write_text("{}\n", encoding="utf-8")

    with pytest.raises(DocumentExtractionError, match="deterministic output"):
        service.parse(ingested.material.id)

    derived = materials.derived_directory(ingested.material.id)
    assert not list(derived.glob(".staging-*"))


@pytest.mark.parametrize(
    ("filename", "payload_factory", "kind", "expected"),
    [
        ("project.docx", _docx_bytes, SourceKind.DOCUMENT, "公开询价"),
        ("project.xlsx", _xlsx_bytes, SourceKind.SHEET, "Project Sheet"),
        ("project.pptx", _pptx_bytes, SourceKind.SLIDE, "PRJ-001"),
    ],
)
def test_office_formats_use_structured_parsers(
    tmp_path: Path,
    filename: str,
    payload_factory,
    kind: SourceKind,
    expected: str,
) -> None:
    materials = _service(tmp_path)
    ingested = materials.ingest_stream(
        BytesIO(payload_factory()),
        original_name=filename,
    )
    parsed = DocumentExtractionService(materials).parse(ingested.material.id)
    assert parsed.document.units
    assert parsed.document.units[0].kind is kind
    combined = "\n".join(
        f"{unit.label}\n{unit.text}" for unit in parsed.document.units
    )
    assert expected in combined
    assert parsed.document.issues == ()


def test_pdf_without_text_layer_requires_ocr_review(tmp_path: Path) -> None:
    materials = _service(tmp_path)
    ingested = materials.ingest_stream(
        BytesIO(_blank_pdf_bytes()),
        original_name="scan.pdf",
        declared_mime_type="application/pdf",
    )
    parsed = DocumentExtractionService(materials).parse(ingested.material.id)
    assert parsed.document.units[0].kind is SourceKind.PAGE
    assert parsed.document.units[0].index == 1
    assert parsed.document.units[0].text == ""
    assert any("no text layer" in issue for issue in parsed.document.issues)
    assert any("OCR/manual review" in issue for issue in parsed.document.issues)


def test_document_character_limit_is_explicit_not_silent(tmp_path: Path) -> None:
    materials = _service(tmp_path)
    ingested = materials.ingest_stream(
        BytesIO(b"abcdefghijklmnopqrstuvwxyz"),
        original_name="long.txt",
        declared_mime_type="text/plain",
    )
    extractor = DocumentExtractor(max_unit_chars=10, max_document_chars=10)
    parsed = DocumentExtractionService(materials, extractor).parse(ingested.material.id)
    assert parsed.document.units[0].text == "abcdefghij"
    assert any("truncated at 10" in issue for issue in parsed.document.issues)


def test_unsupported_image_does_not_fake_text_extraction(tmp_path: Path) -> None:
    materials = _service(tmp_path)
    ingested = materials.ingest_stream(
        BytesIO(b"\x89PNG\r\n\x1a\nREDACTED"),
        original_name="image.png",
        declared_mime_type="image/png",
    )
    with pytest.raises(UnsupportedDocumentType, match="image/png"):
        DocumentExtractionService(materials).parse(ingested.material.id)


def test_field_extraction_success_persists_pending_evidence_and_private_result(
    tmp_path: Path,
) -> None:
    materials = _service(tmp_path)
    ingested = materials.ingest_stream(
        BytesIO(
            "\n".join(
                (
                    "项目编号：PRJ-001",
                    "项目名称：REDACTED PROJECT",
                    "采购方式：公开询价",
                )
            ).encode()
        ),
        original_name="project.txt",
        declared_mime_type="text/plain",
    )
    service = FieldExtractionService(materials, RuleBasedExtractionProvider())

    result = service.extract(ingested.material.id, extraction_id="extract-success")

    assert result.status is ExtractionStatus.SUCCEEDED
    assert result.can_advance is True
    assert result.issues == ()
    assert [proposal.field_name for proposal in result.proposals] == [
        "PR_PrjNo",
        "PR_PrjName",
        "PR_ProcurementMethod",
    ]
    result_file = materials.data_dir / result.result_path
    assert result_file.stat().st_mode & 0o777 == 0o600
    assert json.loads(result_file.read_text(encoding="utf-8"))["id"] == result.id

    stored = materials.storage.get_extraction(result.id)
    assert stored is not None
    assert stored["status"] == "succeeded"
    assert stored["can_advance"] == 1
    fields = stored["fields"]
    assert isinstance(fields, list)
    assert len(fields) == 3
    assert {field["review_status"] for field in fields} == {"pending"}
    assert {field["evidence_valid"] for field in fields} == {1}
    assert {field["source_material_id"] for field in fields} == {
        ingested.material.id
    }
    assert all(field["source_text"] for field in fields)


def test_field_extraction_missing_required_field_stays_in_review(tmp_path: Path) -> None:
    materials = _service(tmp_path)
    ingested = materials.ingest_stream(
        BytesIO("项目编号：PRJ-001\n项目名称：REDACTED PROJECT".encode()),
        original_name="incomplete.txt",
        declared_mime_type="text/plain",
    )

    result = FieldExtractionService(
        materials,
        RuleBasedExtractionProvider(),
    ).extract(ingested.material.id, extraction_id="extract-review")

    assert result.status is ExtractionStatus.NEEDS_REVIEW
    assert result.can_advance is False
    assert any(
        issue.code == "missing_required"
        and issue.field_name == "PR_ProcurementMethod"
        for issue in result.issues
    )
    stored = materials.storage.get_extraction(result.id)
    assert stored is not None
    assert stored["status"] == "needs_review"
    assert stored["issue_count"] >= 1
    assert all(field["review_status"] == "pending" for field in stored["fields"])


def test_field_extraction_provider_failure_is_recorded_without_secret_values(
    tmp_path: Path,
) -> None:
    class FailingProvider:
        name = "failing_test_provider"
        model = "redacted-model"
        version = "1"

        def propose(self, _document, _field_specs):
            raise RuntimeError(
                "provider failed password=DO_NOT_STORE token=DO_NOT_STORE_EITHER"
            )

    materials = _service(tmp_path)
    ingested = materials.ingest_stream(
        BytesIO("项目编号：PRJ-001".encode()),
        original_name="failure.txt",
        declared_mime_type="text/plain",
    )
    service = FieldExtractionService(materials, FailingProvider())

    with pytest.raises(RuntimeError, match="provider failed"):
        service.extract(ingested.material.id, extraction_id="extract-failed")

    stored = materials.storage.get_extraction("extract-failed")
    assert stored is not None
    assert stored["status"] == "failed"
    assert stored["error_type"] == "RuntimeError"
    assert "DO_NOT_STORE" not in stored["error_message"]
    assert stored["error_message"].count("<redacted>") == 2
    assert stored["fields"] == []
    result_path = (
        materials.derived_directory(ingested.material.id)
        / "extractions"
        / "extract-failed"
        / "result.json"
    )
    assert not result_path.exists()


def test_extraction_cli_runs_offline_through_the_same_service(
    tmp_path: Path,
    capsys: pytest.CaptureFixture[str],
) -> None:
    data_dir = tmp_path / "data"
    materials = MaterialService(data_dir=data_dir)
    ingested = materials.ingest_stream(
        BytesIO(
            "项目编号：PRJ-001\n项目名称：REDACTED PROJECT\n采购方式：公开询价".encode()
        ),
        original_name="cli-project.txt",
        declared_mime_type="text/plain",
    )

    exit_code = extract_cli.main(
        [
            ingested.material.id,
            "--data-dir",
            str(data_dir),
            "--json",
        ]
    )

    assert exit_code == 0
    payload = json.loads(capsys.readouterr().out)
    assert payload["extraction"]["status"] == "succeeded"
    assert payload["extraction"]["can_advance"] is True
    extraction_id = payload["extraction"]["id"]
    stored = materials.storage.get_extraction(extraction_id)
    assert stored is not None
    assert stored["status"] == "succeeded"


def test_extraction_cli_rejects_invalid_threshold_without_running(
    tmp_path: Path,
    capsys: pytest.CaptureFixture[str],
) -> None:
    assert (
        extract_cli.main(
            [
                "missing",
                "--data-dir",
                str(tmp_path),
                "--confidence-threshold",
                "1.1",
            ]
        )
        == 2
    )
    assert "between 0 and 1" in capsys.readouterr().err
    assert not list(tmp_path.iterdir())
